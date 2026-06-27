"""
Análise Executiva — Export PPTX nativo (P1).

export_to_pptx_bytes(deck_spec) -> bytes  (assinatura espelha
email_service.export_to_excel_bytes: recebe dict, devolve bytes via BytesIO).

Princípios:
- Gráficos NATIVOS do PowerPoint (add_chart) alimentados pelo data{rows} — sem
  screenshot de Chart.js, sem matplotlib/kaleido. O .pptx sai editável.
- Marca Claro (vermelho #E30613), tag de seção, número-herói grande, caixa de
  "Ação recomendada", rodapé de fonte + nº do slide.
- O SQL de cada insight vai para as speaker notes (auditabilidade fora da app).
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ---- Marca -----------------------------------------------------------------
RED   = RGBColor(0xE3, 0x06, 0x13)
INK   = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
LINE  = RGBColor(0xE2, 0xE5, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)

SW, SH = 13.333, 7.5
M = 0.6
_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}


def _short(text, n: int = 320) -> str:
    """Limpa markdown (**, crases) e trunca texto para caber numa caixa de slide
    — defensivo contra narrativas longas que transbordam e sobrepõem o número-
    herói e as ações no PPTX."""
    if not text:
        return ""
    import re as _re
    t = str(text).replace("**", "").replace("`", "")
    t = _re.sub(r"\s+", " ", t).strip()
    if len(t) <= n:
        return t
    cut = t[:n]
    last = max(cut.rfind(". "), cut.rfind("? "), cut.rfind("! "))
    if last > 60:
        return cut[:last + 1]
    sp = cut.rfind(" ")
    return (cut[:sp] if sp > 60 else cut).rstrip() + "…"


def _clean_md_lines(text) -> list[str]:
    """Limpa markdown (**, crases, # ) e quebra a narrativa em linhas/bullets legíveis,
    PRESERVANDO o conteúdo ÍNTEGRO (sem truncar). Bullets inline ' - ' e numeração
    'N. ' viram quebras de linha; mantém parágrafos. Usado no layout texto-forward."""
    if not text:
        return []
    import re as _re
    t = str(text).replace("\r\n", "\n").replace("\r", "\n")
    t = t.replace("**", "").replace("`", "")
    t = _re.sub(r"(?m)^\s{0,3}#{1,6}\s*", "", t)               # cabeçalhos markdown
    t = _re.sub(r"\s+-\s+", "\n• ", t)                          # bullets inline ' - '
    t = _re.sub(r"(?<=[\s.;:])(\d{1,2})\.\s+", r"\n\1. ", t)    # numeração 'N. '
    out: list[str] = []
    for ln in t.split("\n"):
        ln = _re.sub(r"[ \t]+", " ", ln).strip()
        ln = _re.sub(r"^[*\-]\s+", "• ", ln)                    # bullet no início da linha
        if ln:
            out.append(ln)
    return out


def _para(slide, l, t, w, h, lines, size=12, color=INK, name="Calibri"):
    """Caixa multi-parágrafo com autoajuste (SHRINK_TEXT_ON_OVERFLOW): o PowerPoint
    reduz a fonte para caber TODO o conteúdo — mostra a narrativa íntegra sem truncar."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.SHRINK_TEXT_ON_OVERFLOW
    except Exception:
        pass
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    items = lines if isinstance(lines, (list, tuple)) else [str(lines)]
    first = True
    for ln in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = str(ln)
        r.font.size = Pt(size)
        r.font.name = name
        r.font.color.rgb = color
    return tb


def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    try:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
    except Exception:
        pass
    return s


def _txt(slide, l, t, w, h, text, size=12, bold=False, color=INK, align="l",
         name="Calibri", italic=False, anchor=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = _ALIGN.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = str(text if text is not None else "")
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = name
    r.font.color.rgb = color
    return tb


def _bullets(slide, l, t, w, h, items, size=12, color=INK, name="Calibri",
             bullet="•  ", space_after=6):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    first = True
    for it in items or []:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = bullet + str(it)
        r.font.size = Pt(size)
        r.font.name = name
        r.font.color.rgb = color
    return tb


def _rect(slide, l, t, w, h, fill=LIGHT, line=None, rounded=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _section_tag(slide, text):
    _txt(slide, M, 0.36, 8, 0.3, (text or "").upper(), size=12, bold=True, color=RED)


def _footer(slide, text, page):
    _txt(slide, M, SH - 0.45, SW - 2 * M - 0.6, 0.3, text or "", size=8.5, color=MUTED)
    _txt(slide, SW - M - 0.6, SH - 0.45, 0.6, 0.3, str(page), size=8.5, color=MUTED, align="r")


def _conf_seal(slide, conf):
    if not conf or not conf.get("level"):
        return
    level = conf.get("level")
    col = GREEN if level == "Alta" else (RED if level == "Baixa" else AMBER)
    _txt(slide, SW - M - 2.6, 0.34, 2.6, 0.3, f"Confiança: {level}", size=10, bold=True, color=col, align="r")


def _pct_display(value, obj):
    """Garante o símbolo "%" em valores percentuais no PPTX (espelha _pctDisplay
    do front), p/ decks gerados ANTES do guard de geração. Recomputa do value_raw
    quando o fmt é percentual ou quando rótulo/coluna indicam percentual e o valor
    está numa faixa plausível; senão devolve a string original inalterada."""
    s = "" if value is None else str(value)
    if "%" in s:
        return s
    obj = obj or {}
    fmt = obj.get("fmt") or ""
    vr = obj.get("value_raw")
    try:
        from app.services.report_service import _format_value, _coerce_number
        from app.services.exec_analysis_service import _looks_percent
    except Exception:
        return s
    if fmt.startswith("percent"):
        return _format_value(vr, fmt) if vr is not None else s
    if vr is None or not _looks_percent(obj.get("label"), obj.get("column"), obj.get("caption")):
        return s
    n = _coerce_number(vr)
    if n is None:
        return s
    if -1.0 <= n <= 1.0:
        return _format_value(vr, "percent_2")
    if abs(n) <= 150.0:
        return _format_value(vr, "percent_raw")
    return s


# ---------------------------------------------------------------------------
# Builders por tipo de slide
# ---------------------------------------------------------------------------

def _cover(prs, slide_spec, deck):
    s = _slide(prs)
    _rect(s, 0, 0, 0.22, SH, fill=RED, rounded=False)
    _txt(s, 0.9, 2.3, 11.5, 2.0, slide_spec.get("title") or deck.get("title", ""),
         size=40, bold=True, color=INK)
    sub = slide_spec.get("subtitle") or deck.get("thesis", "")
    if sub:
        _txt(s, 0.9, 4.4, 10.8, 1.6, sub, size=18, color=MUTED)
    _txt(s, 0.9, 6.6, 9, 0.4, slide_spec.get("date_label") or deck.get("date_label", ""),
         size=11, bold=True, color=RED)


def _header(s, slide_spec):
    _section_tag(s, slide_spec.get("section", ""))
    _txt(s, M, 0.68, SW - 2 * M, 0.7, slide_spec.get("title", ""), size=26, bold=True, color=INK)


def _sintese(prs, sp, deck, page):
    s = _slide(prs)
    _header(s, sp)
    if sp.get("thesis"):
        _txt(s, M, 1.45, SW - 2 * M, 0.8, sp["thesis"], size=15, bold=True, color=INK)
    callouts = sp.get("callouts") or []
    if callouts:
        n = len(callouts)
        gap = 0.2
        cw = (SW - 2 * M - gap * (n - 1)) / n
        for i, c in enumerate(callouts):
            x = M + i * (cw + gap)
            _rect(s, x, 2.45, cw, 1.25, fill=LIGHT)
            _txt(s, x + 0.1, 2.58, cw - 0.2, 0.7, _pct_display(c.get("value", ""), c), size=30, bold=True, color=RED, align="c")
            _txt(s, x + 0.1, 3.28, cw - 0.2, 0.35, c.get("label", ""), size=10, color=MUTED, align="c")
    colw = (SW - 2 * M - 0.4) / 2
    y = 4.1
    _txt(s, M, y, colw, 0.35, "O que os dados mostram", size=13, bold=True, color=INK)
    _bullets(s, M, y + 0.45, colw, 2.2, sp.get("o_que_mostram") or [], size=12, color=INK)
    x2 = M + colw + 0.4
    _txt(s, x2, y, colw, 0.35, "Implicação estratégica", size=13, bold=True, color=RED)
    _bullets(s, x2, y + 0.45, colw, 2.2, sp.get("implicacao") or [], size=12, color=INK)
    _footer(s, deck.get("source_footer", ""), page)


def _sowhat(prs, sp, deck, page):
    s = _slide(prs)
    _header(s, sp)
    pilares = sp.get("pilares") or []
    n = max(1, len(pilares))
    gap = 0.3
    cw = (SW - 2 * M - gap * (n - 1)) / n
    for i, p in enumerate(pilares):
        x = M + i * (cw + gap)
        _rect(s, x, 1.7, cw, 2.7, fill=LIGHT)
        _txt(s, x + 0.25, 1.85, 0.8, 0.7, str(p.get("n", i + 1)), size=34, bold=True, color=RED)
        _txt(s, x + 0.25, 2.7, cw - 0.5, 0.5, p.get("title", ""), size=15, bold=True, color=INK)
        _txt(s, x + 0.25, 3.2, cw - 0.5, 1.1, p.get("text", ""), size=11.5, color=MUTED)
    if sp.get("thesis"):
        _rect(s, M, 4.8, SW - 2 * M, 1.0, fill=RGBColor(0xFC, 0xE9, 0xEA))
        _txt(s, M + 0.25, 5.0, SW - 2 * M - 0.5, 0.7, "Tese: " + sp["thesis"], size=13, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, deck.get("source_footer", ""), page)


def _add_chart(slide, chart, l, t, w, h):
    try:
        labels = [str(x) for x in (chart.get("labels") or [])]
        values = [float(v) if v is not None else 0.0 for v in (chart.get("values") or [])]
        if len(labels) < 2 or len(values) < 2:
            return False
        ctype_name = (chart.get("type") or "bar").lower()
        is_pie = ctype_name in ("pie", "doughnut")
        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series(chart.get("y_field", "valor"), values)
        ctype = {
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "area": XL_CHART_TYPE.AREA,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }.get(ctype_name, XL_CHART_TYPE.COLUMN_CLUSTERED)
        gframe = slide.shapes.add_chart(ctype, Inches(l), Inches(t), Inches(w), Inches(h), cd)
        ch = gframe.chart
        ch.has_legend = is_pie
        try:
            ch.has_title = False
            plot = ch.plots[0]
            plot.has_data_labels = False
            if not is_pie:
                ser = plot.series[0]
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = RED
        except Exception:
            pass
        return True
    except Exception:
        return False


def _insight_notes(slide, sp):
    """SQL + fundamento da confiança nas speaker notes (auditabilidade; PPTX não tem
    tooltip, então o critério/motivo da confiança vai para as notas)."""
    note_parts = []
    _cf = sp.get("confidence") or {}
    if _cf.get("level"):
        _crit = ("Criterios — Alta: fonte catalogada, completude >=90% e qualidade "
                 ">=70%; Media: sem catalogo ou completude 70-90%/qualidade <70%; "
                 "Baixa: completude <70%, 0 linhas ou sem numero-heroi.")
        note_parts.append(f"Confianca: {_cf.get('level')}"
                          + (f" — {_cf.get('reason')}" if _cf.get("reason") else "")
                          + "\n" + _crit)
    if sp.get("sql"):
        note_parts.append("SQL do número-chave:\n" + sp["sql"])
    if note_parts:
        try:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(note_parts)
        except Exception:
            pass


def _insight(prs, sp, deck, page):
    s = _slide(prs)
    _header(s, sp)
    _conf_seal(s, sp.get("confidence"))
    if sp.get("subtitle"):
        _txt(s, M, 1.32, SW - 2 * M, 0.4, sp["subtitle"], size=13, color=MUTED, italic=True)
    has_chart = bool(sp.get("chart"))
    hero = sp.get("hero") or {}
    actions = sp.get("actions") or []
    cz = sp.get("causal")
    foot = _src_footer_for(sp) or deck.get("source_footer", "")
    narr_lines = _clean_md_lines(sp.get("narrative"))
    narr_len = sum(len(x) for x in narr_lines)
    sub_off = 0.45 if sp.get("subtitle") else 0.0

    # --- Slides "aprofundamento": narrativa longa SEM gráfico → layout texto-forward
    #     que mostra a narrativa ÍNTEGRA (autoajuste), com o número-chave compacto no
    #     topo — em vez de truncar para caber num cantinho. ---
    if narr_lines and narr_len > 360 and not has_chart:
        y0 = 1.5 + sub_off
        hv = _pct_display(hero.get("value", "—"), hero)
        if hv and str(hv) != "—":
            _txt(s, M, y0, 3.1, 0.8, hv, size=28, bold=True, color=RED, anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, M + 3.25, y0 + 0.05, SW - 2 * M - 3.25, 0.4, hero.get("label", ""), size=13, bold=True, color=INK)
            if hero.get("caption"):
                _txt(s, M + 3.25, y0 + 0.45, SW - 2 * M - 3.25, 0.4, hero["caption"], size=10, color=MUTED)
            ny = y0 + 1.0
        else:
            ny = y0  # sem número (ex.: "Sem dados") → narrativa ocupa desde o topo
        ab = SH - 0.55                       # base acima do rodapé
        if actions:
            _txt(s, M, ab - 0.9, SW - 2 * M, 0.3, "Ação recomendada", size=12, bold=True, color=RED)
            _bullets(s, M, ab - 0.58, SW - 2 * M, 0.5, actions, size=11, color=INK, space_after=3)
        nh = (ab - 1.0 if actions else ab) - ny
        _para(s, M, ny, SW - 2 * M, max(1.0, nh), narr_lines, size=12, color=INK)
        _footer(s, foot, page)
        _insight_notes(s, sp)
        return

    # --- Layout padrão: número-herói grande + narrativa curta/limpa acima ---
    left_w = 5.6 if has_chart else (SW - 2 * M)
    if sp.get("narrative"):
        _txt(s, M, 1.85, left_w, 1.05, _short(sp["narrative"], 240 if has_chart else 360), size=11, color=INK)
    _txt(s, M, 3.0, left_w, 1.1, _pct_display(hero.get("value", "—"), hero), size=52, bold=True, color=RED)
    _txt(s, M, 4.15, left_w, 0.35, hero.get("label", ""), size=13, bold=True, color=INK)
    if hero.get("caption"):
        _txt(s, M, 4.5, left_w, 0.35, hero["caption"], size=11, color=MUTED)
    if cz:
        ccol = GREEN if cz.get("significant") else AMBER
        _txt(s, M, 4.85, left_w, 0.5,
             f"Efeito causal (PSM): {cz.get('effect_label','')}  ·  {cz.get('caveat','')}",
             size=9.5, bold=True, color=ccol)
    if has_chart:
        _add_chart(s, sp["chart"], 6.5, 1.7, SW - M - 6.5, 3.4)
    if actions:
        ay = 5.5 if cz else 5.25
        _txt(s, M, ay, left_w, 0.3, "Ação recomendada", size=12, bold=True, color=RED)
        _bullets(s, M, ay + 0.35, left_w, 1.0, actions, size=11, color=INK, space_after=3)
    _footer(s, foot, page)
    _insight_notes(s, sp)


def _src_footer_for(sp):
    src = sp.get("source") or {}
    tables = ", ".join(src.get("tables", [])[:3]) if src else ""
    if not tables:
        return ""
    out = "Fonte: " + tables
    comp = src.get("min_completeness")
    if isinstance(comp, (int, float)):
        out += f" | completude {int(comp)}%"
    return out + " | Análise interna"


def _estrategia(prs, sp, deck, page):
    s = _slide(prs)
    _header(s, sp)
    frentes = sp.get("frentes") or []
    cols, gap = 2, 0.3
    cw = (SW - 2 * M - gap) / cols
    rh = 2.1
    for i, f in enumerate(frentes[:4]):
        row, col = divmod(i, cols)
        x = M + col * (cw + gap)
        y = 1.7 + row * (rh + 0.25)
        _rect(s, x, y, cw, rh, fill=LIGHT)
        _txt(s, x + 0.25, y + 0.18, 0.8, 0.6, str(f.get("n", i + 1)), size=30, bold=True, color=RED)
        _txt(s, x + 1.1, y + 0.22, cw - 1.3, 0.6, f.get("title", ""), size=15, bold=True, color=INK)
        _txt(s, x + 0.25, y + 0.95, cw - 0.5, 1.0, f.get("text", ""), size=11.5, color=MUTED)
    _footer(s, deck.get("source_footer", ""), page)


def _roadmap(prs, sp, deck, page):
    s = _slide(prs)
    _header(s, sp)
    sprints = sp.get("sprints") or []
    n = max(1, len(sprints))
    gap = 0.2
    cw = (SW - 2 * M - gap * (n - 1)) / n
    for i, sp_ in enumerate(sprints[:4]):
        x = M + i * (cw + gap)
        _rect(s, x, 1.7, cw, 0.55, fill=RED)
        _txt(s, x + 0.1, 1.78, cw - 0.2, 0.4, sp_.get("period", ""), size=12, bold=True, color=WHITE, align="c")
        _rect(s, x, 2.3, cw, 3.4, fill=LIGHT)
        _txt(s, x + 0.18, 2.45, cw - 0.36, 0.7, sp_.get("title", ""), size=13, bold=True, color=INK)
        _bullets(s, x + 0.18, 3.2, cw - 0.36, 2.4, sp_.get("bullets") or [], size=10.5, color=MUTED, space_after=4)
    _footer(s, deck.get("source_footer", ""), page)


def _kpis(prs, sp, deck, page):
    s = _slide(prs)
    _header(s, sp)
    metas = sp.get("metas") or []
    bottom_y = 3.3
    if metas:
        n = len(metas)
        gap = 0.2
        cw = (SW - 2 * M - gap * (n - 1)) / n
        card_top, card_h = 1.55, 2.35
        for i, m in enumerate(metas):
            x = M + i * (cw + gap)
            _rect(s, x, card_top, cw, card_h, fill=LIGHT)
            val = str(m.get("value", "") or "")
            # Fonte adaptativa ao tamanho do texto: valores longos usam um corpo
            # menor para caber sem sobrepor o rótulo — sem cortar conteúdo.
            L = len(val)
            vsize = (30 if L <= 6 else 24 if L <= 12 else 19 if L <= 20
                     else 15 if L <= 32 else 13 if L <= 48 else 11)
            _txt(s, x + 0.12, card_top + 0.12, cw - 0.24, 1.45, val,
                 size=vsize, bold=True, color=RED, align="c", anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, x + 0.12, card_top + 1.62, cw - 0.24, card_h - 1.74,
                 m.get("label", ""), size=10, color=MUTED, align="c", anchor=MSO_ANCHOR.TOP)
        bottom_y = card_top + card_h + 0.35  # empurra a seção de baixo p/ não colidir
    colw = (SW - 2 * M - 0.4) / 2
    y = bottom_y
    _txt(s, M, y, colw, 0.35, "Ritual de gestão", size=13, bold=True, color=INK)
    _bullets(s, M, y + 0.45, colw, 2.2, sp.get("ritual") or [], size=12, color=INK)
    x2 = M + colw + 0.4
    _txt(s, x2, y, colw, 0.35, "Papéis recomendados", size=13, bold=True, color=RED)
    donos = []
    for d in sp.get("donos") or []:
        if isinstance(d, dict):
            donos.append(f"{d.get('role','')}: {d.get('scope','')}")
        else:
            donos.append(str(d))
    _bullets(s, x2, y + 0.45, colw, 2.2, donos, size=12, color=INK)
    _footer(s, deck.get("source_footer", ""), page)


_BUILDERS = {
    "sintese": _sintese, "sowhat": _sowhat, "insight": _insight,
    "estrategia": _estrategia, "roadmap": _roadmap, "kpis": _kpis,
}


def export_to_pptx_bytes(deck: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    page = 0
    for sp in deck.get("slides", []):
        page += 1
        stype = sp.get("type")
        try:
            if stype == "cover":
                _cover(prs, sp, deck)
            elif stype in _BUILDERS:
                _BUILDERS[stype](prs, sp, deck, page)
            else:
                # tipo desconhecido — slide de texto simples (não quebra o export)
                s = _slide(prs)
                _header(s, sp)
                _footer(s, deck.get("source_footer", ""), page)
        except Exception:
            # um slide problemático não pode derrubar o deck inteiro
            try:
                s = _slide(prs)
                _txt(s, M, 0.7, SW - 2 * M, 0.7, sp.get("title", "Slide"), size=24, bold=True, color=INK)
            except Exception:
                pass
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
